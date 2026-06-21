import pptx

prs = pptx.Presentation('node_framework.pptx')

for slide in prs.slides:
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        
        # We can check the whole shape.text first to see if it needs replacement.
        # This helps in case the word is split across runs.
        full_text = shape.text
        needs_update = False
        if "Progress Report" in full_text or "Recent Work Progress" in full_text or "Critical Bugs Fixed" in full_text or "Next Steps" in full_text:
            needs_update = True
            
        if needs_update:
            # We will clear the text frame and insert the replaced text, keeping the first run's formatting
            new_text = full_text.replace("Progress Report", "Final Report")
            new_text = new_text.replace("Recent Work Progress", "Completed Features")
            new_text = new_text.replace("Critical Bugs Fixed", "Key Challenges Resolved")
            new_text = new_text.replace("Next Steps", "Future Scope")
            
            # Save format of the first run
            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                font = shape.text_frame.paragraphs[0].runs[0].font
            else:
                font = None
                
            shape.text = new_text
            
            if font and shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                new_run = shape.text_frame.paragraphs[0].runs[0]
                new_run.font.name = font.name
                new_run.font.size = font.size
                new_run.font.bold = font.bold
                new_run.font.italic = font.italic
                new_run.font.color.rgb = font.color.rgb if hasattr(font.color, 'rgb') and font.color.type == 1 else None

        # Check footers or smaller things
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if '2025' in run.text:
                    run.text = run.text.replace('2025', '2026')
                if 'WebNode Framework  •  Progress Report' in run.text:
                    run.text = run.text.replace('Progress Report', 'Final Report')

# Let's fix slide 9's Future Scope items
slide9 = prs.slides[8]
for shape in slide9.shapes:
    if hasattr(shape, 'text_frame'):
        if 'Visual Node Editor UI (drag-and-drop graph builder)' in shape.text:
            shape.text = (
                "Implemented Visual Node Editor UI (drag-and-drop builder)\n"
                "Auto-regeneration of main.py via Deploy\n"
                "SQLite with WAL mode for Thread-Safe persistence\n"
                "Global Fetch Override for CSRF Tokens in API\n"
                "Future: CacheNode, WebSocketNode, AuthNode expansions"
            )

prs.save('node_framework.pptx')
print("Done updating PPTX")
